import os
import sys
import io

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..")))

import docx
import pptx
from fastapi.testclient import TestClient
from backend.app.main import app
from backend.app.core.security import create_access_token
from backend.app.services.mongo_service import MongoService

client = TestClient(app)

def create_sample_docx_bytes() -> bytes:
    doc = docx.Document()
    doc.add_heading("Q3 Performance Overview", level=1)
    doc.add_paragraph("This section summarizes overall financial and digital metrics.")
    doc.add_heading("Detailed Breakdown", level=2)
    doc.add_paragraph("Key metric targets achieved.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()

def create_sample_pptx_bytes() -> bytes:
    prs = pptx.Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(0, 0, 100, 100)
    tf = txBox.text_frame
    tf.text = "Executive Summary Slide Content"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def get_auth_headers(client_id: str = "client_test1"):
    token = create_access_token(client_id=client_id)
    return {"Authorization": f"Bearer {token}"}

def test_template_upload_and_lifecycle():
    client_id = "client_test1"
    headers = get_auth_headers(client_id)
    docx_bytes = create_sample_docx_bytes()

    # 1. Upload DOCX template
    response = client.post(
        "/api/v1/templates/upload",
        data={
            "client_id": client_id,
            "template_name": "Test Quarterly DOCX",
            "description": "Used for quarterly executive briefings."
        },
        files={"file": ("report_template.docx", docx_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
        headers=headers
    )
    assert response.status_code == 200, response.text
    data = response.json()
    assert data["status"] == "success"
    template_id = data["template_id"]
    assert template_id.startswith("tmpl_")

    # 2. List templates for client
    list_res = client.get(f"/api/v1/templates/list?client_id={client_id}", headers=headers)
    assert list_res.status_code == 200
    list_data = list_res.json()
    assert list_data["status"] == "success"
    assert len(list_data["templates"]) >= 1
    found = [t for t in list_data["templates"] if t["template_id"] == template_id]
    assert len(found) == 1
    assert found[0]["template_name"] == "Test Quarterly DOCX"
    assert found[0]["file_type"] == "docx"

    # 3. Download template
    dl_res = client.get(f"/api/v1/templates/download/{template_id}", headers=headers)
    assert dl_res.status_code == 200
    assert len(dl_res.content) == len(docx_bytes)

    # 4. Orchestrate using template_id
    orch_res = client.post(
        "/api/v1/orchestrate",
        json={
            "client_id": client_id,
            "user_prompt": "Generate report with uploaded docx",
            "template_id": template_id
        },
        headers=headers
    )
    assert orch_res.status_code == 200
    orch_data = orch_res.json()
    assert orch_data["status"] == "success"
    assert orch_data["tiptap_json"]["type"] == "doc"

def test_template_validation_and_security():
    headers = get_auth_headers("client_test1")

    # Invalid extension (.txt)
    txt_bytes = b"Hello world text file"
    res1 = client.post(
        "/api/v1/templates/upload",
        data={"client_id": "client_test1", "template_name": "Invalid File", "description": "Invalid format test"},
        files={"file": ("invalid.txt", txt_bytes, "text/plain")},
        headers=headers
    )
    assert res1.status_code == 400
    assert "Only .docx and .pptx files are allowed" in res1.text

    # File size over limit (16MB)
    big_bytes = b"0" * (16 * 1024 * 1024)
    res2 = client.post(
        "/api/v1/templates/upload",
        data={"client_id": "client_test1", "template_name": "Huge File", "description": "Too large"},
        files={"file": ("huge.docx", big_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
        headers=headers
    )
    assert res2.status_code == 400
    assert "File size exceeds maximum allowed limit of 15MB" in res2.text

    # Tenant mismatch (auth client_test1 trying to upload for client_other)
    res3 = client.post(
        "/api/v1/templates/upload",
        data={"client_id": "client_other", "template_name": "Mismatch", "description": "Tenant mismatch test"},
        files={"file": ("sample.docx", create_sample_docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
        headers=headers
    )
    assert res3.status_code == 403 or "Tenant isolation mismatch" in res3.text

def test_pptx_upload():
    client_id = "client_test1"
    headers = get_auth_headers(client_id)
    pptx_bytes = create_sample_pptx_bytes()

    response = client.post(
        "/api/v1/templates/upload",
        data={
            "client_id": client_id,
            "template_name": "Test Slide Deck PPTX",
            "description": "Used for client presentation decks."
        },
        files={"file": ("deck.pptx", pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        headers=headers
    )
    assert response.status_code == 200
    data = response.json()
    assert data["status"] == "success"
    assert data["template_name"] == "Test Slide Deck PPTX"
