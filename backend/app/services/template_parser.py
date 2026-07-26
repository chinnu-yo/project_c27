import io
from typing import Dict, Any, List
import docx
import pptx

def parse_docx_structure(file_bytes: bytes) -> Dict[str, Any]:
    """
    Parses a .docx file buffer using python-docx.
    Extracts heading hierarchy, paragraph styles, section ordering, tables, and constructs
    a Tiptap-compatible schema blueprint for report orchestration.
    """
    doc = docx.Document(io.BytesIO(file_bytes))
    
    sections_info: List[Dict[str, Any]] = []
    for idx, sec in enumerate(doc.sections, start=1):
        orientation = "LANDSCAPE" if getattr(sec, "orientation", None) == docx.enum.section.WD_ORIENT.LANDSCAPE else "PORTRAIT"
        sections_info.append({
            "section_number": idx,
            "orientation": orientation,
            "page_width_inches": round(sec.page_width.inches, 2) if sec.page_width else None,
            "page_height_inches": round(sec.page_height.inches, 2) if sec.page_height else None
        })

    headings_info: List[Dict[str, Any]] = []
    paragraphs_info: List[Dict[str, Any]] = []
    tables_info: List[Dict[str, Any]] = []
    tiptap_nodes: List[Dict[str, Any]] = []
    text_layout_lines: List[str] = []

    for p in doc.paragraphs:
        text = p.text.strip()
        style_name = p.style.name if p.style else "Normal"
        
        if not text:
            continue

        # Check if style indicates a heading
        if style_name.startswith("Heading"):
            try:
                level = int(style_name.replace("Heading", "").strip())
            except ValueError:
                level = 1
            headings_info.append({
                "level": level,
                "text": text,
                "style_name": style_name
            })
            tiptap_nodes.append({
                "type": "heading",
                "attrs": {"level": min(max(level, 1), 6)},
                "content": [{"type": "text", "text": text}]
            })
            text_layout_lines.append(f"{'#' * level} {text}")
        else:
            paragraphs_info.append({
                "style_name": style_name,
                "text": text
            })
            tiptap_nodes.append({
                "type": "paragraph",
                "content": [{"type": "text", "text": text}]
            })
            text_layout_lines.append(text)

    # Extract tables from .docx document
    for t_idx, table in enumerate(doc.tables, start=1):
        table_rows_data: List[List[str]] = []
        tiptap_row_nodes: List[Dict[str, Any]] = []

        for r_idx, row in enumerate(table.rows):
            row_cells = [cell.text.strip() for cell in row.cells]
            table_rows_data.append(row_cells)

            cell_type = "tableHeader" if r_idx == 0 else "tableCell"
            tiptap_cells = []
            for cell_text in row_cells:
                tiptap_cells.append({
                    "type": cell_type,
                    "content": [{
                        "type": "paragraph",
                        "content": [{"type": "text", "text": cell_text or " "}]
                    }]
                })
            tiptap_row_nodes.append({
                "type": "tableRow",
                "content": tiptap_cells
            })

        tables_info.append({
            "table_number": t_idx,
            "row_count": len(table.rows),
            "col_count": len(table.columns) if table.columns else 0,
            "rows": table_rows_data
        })

        if tiptap_row_nodes:
            tiptap_nodes.append({
                "type": "table",
                "content": tiptap_row_nodes
            })
            formatted_rows = [" | ".join(r) for r in table_rows_data]
            text_layout_lines.append(f"\n[TABLE FORMAT:\n" + "\n".join(formatted_rows) + "\n]")

    raw_text_layout = "\n".join(text_layout_lines) if text_layout_lines else (
        "1. Executive Financial Summary\n"
        "2. Invoicing & Payment Breakdown\n"
        "3. CRM Pipeline\n"
        "4. Financial Recommendations"
    )

    # Ensure a fallback tiptap_schema_blueprint node exists
    if not tiptap_nodes:
        tiptap_nodes = [
            {"type": "heading", "attrs": {"level": 1}, "content": [{"type": "text", "text": "Extracted Document Template"}]},
            {"type": "paragraph", "content": [{"type": "text", "text": "Template section baseline."}]}
        ]

    return {
        "file_type": "docx",
        "sections": sections_info,
        "headings": headings_info,
        "paragraphs": paragraphs_info,
        "tables": tables_info,
        "raw_text_layout": raw_text_layout,
        "template_text_content": raw_text_layout,
        "structure_summary": {
            "section_count": len(sections_info),
            "heading_count": len(headings_info),
            "paragraph_count": len(paragraphs_info),
            "table_count": len(tables_info)
        },
        "tiptap_schema_blueprint": {
            "type": "doc",
            "content": tiptap_nodes
        }
    }


def parse_pptx_structure(file_bytes: bytes) -> Dict[str, Any]:
    """
    Parses a .pptx file buffer using python-pptx.
    Extracts slide layout names, placeholder structure, text content, and constructs
    a Tiptap-compatible schema blueprint for report orchestration.
    """
    prs = pptx.Presentation(io.BytesIO(file_bytes))
    
    slides_info: List[Dict[str, Any]] = []
    tiptap_nodes: List[Dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name if slide.slide_layout else f"Layout {idx}"
        placeholders: List[Dict[str, Any]] = []
        
        slide_heading_text = f"Slide {idx}: {layout_name}"
        tiptap_nodes.append({
            "type": "heading",
            "attrs": {"level": 2},
            "content": [{"type": "text", "text": slide_heading_text}]
        })

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                txt = shape.text.strip()
                shape_name = shape.name
                p_type = str(shape.shape_type)
                
                placeholders.append({
                    "name": shape_name,
                    "type": p_type,
                    "text": txt
                })

                tiptap_nodes.append({
                    "type": "paragraph",
                    "content": [{"type": "text", "text": f"[{shape_name}] {txt}"}]
                })

        slides_info.append({
            "slide_number": idx,
            "layout_name": layout_name,
            "placeholders": placeholders
        })

    if not tiptap_nodes:
        tiptap_nodes = [
            {"type": "heading", "attrs": {"level": 1}, "content": [{"type": "text", "text": "Extracted Presentation Template"}]},
            {"type": "paragraph", "content": [{"type": "text", "text": "Slide deck layout baseline."}]}
        ]

    return {
        "file_type": "pptx",
        "slides": slides_info,
        "structure_summary": {
            "slide_count": len(slides_info)
        },
        "tiptap_schema_blueprint": {
            "type": "doc",
            "content": tiptap_nodes
        }
    }


def parse_template_file(file_bytes: bytes, filename: str) -> Dict[str, Any]:
    """
    Main parser entrypoint that inspects file extension and extracts structure.
    """
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext == "docx":
        return parse_docx_structure(file_bytes)
    elif ext == "pptx":
        return parse_pptx_structure(file_bytes)
    else:
        raise ValueError(f"Unsupported file extension: .{ext}")
